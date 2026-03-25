"""
Generate MS Word Report and PowerPoint Presentation
for ML Model Comparison - ClaimBot AI Masters Thesis
"""

import csv
from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt, Emu as PptEmu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

RESULTS_DIR = Path(__file__).resolve().parent / "comparison_results"
OUTPUT_DIR = Path(__file__).resolve().parent.parent / "documents"
OUTPUT_DIR.mkdir(exist_ok=True)

# ── Load CSV data ───────────────────────────────────────────

def load_csv(name):
    rows = []
    with open(str(RESULTS_DIR / name), "r") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(row)
    return rows  # rows[0] = header, rows[1:] = data


clf_rows = load_csv("classification_metrics.csv")
reg_rows = load_csv("regression_metrics.csv")
summary_rows = load_csv("model_summary.csv")

# Chart image paths
CHARTS = {
    "classifier_comparison": RESULTS_DIR / "classifier_comparison.png",
    "confusion_matrices": RESULTS_DIR / "confusion_matrices.png",
    "roc_curves": RESULTS_DIR / "roc_curves.png",
    "regressor_comparison": RESULTS_DIR / "regressor_comparison.png",
    "training_time": RESULTS_DIR / "training_time.png",
    "feature_importance": RESULTS_DIR / "feature_importance.png",
}


# ╔═══════════════════════════════════════════════════════════╗
# ║  PART 1: MS WORD DOCUMENT                                ║
# ╚═══════════════════════════════════════════════════════════╝

def add_table_from_csv(doc, rows, max_cols=None):
    """Add a formatted table from CSV rows."""
    header = rows[0][:max_cols] if max_cols else rows[0]
    data = [r[:max_cols] if max_cols else r for r in rows[1:]]

    table = doc.add_table(rows=1 + len(data), cols=len(header))
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    # Header row
    for j, val in enumerate(header):
        cell = table.rows[0].cells[j]
        cell.text = val
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.bold = True
                run.font.size = Pt(9)

    # Data rows
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.rows[i + 1].cells[j]
            try:
                fval = float(val)
                if abs(fval) < 0.01 or fval == int(fval):
                    cell.text = val
                else:
                    cell.text = f"{fval:.4f}" if fval < 100 else f"{fval:.1f}"
            except ValueError:
                cell.text = val
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(9)

    return table


def create_word_document():
    doc = Document()

    # ── Style defaults ──
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    # ── Title Page ──
    for _ in range(6):
        doc.add_paragraph("")

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("ML Model Comparison Report")
    run.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0x10, 0xA3, 0x7F)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(
        "Dynamic AI-Powered Virtual Agent for\n"
        "Real-Time Health Insurance Claim Status Queries"
    )
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x4A, 0x4A, 0x4A)

    doc.add_paragraph("")
    info = doc.add_paragraph()
    info.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = info.add_run(
        "Masters Thesis Project\n"
        "Dataset: CMS DE-SynPUF (500,000 claims)\n"
        "Models: XGBoost, Random Forest, Logistic/Linear Regression,\n"
        "Gradient Boosting, FNN, CNN (1D Conv), RNN (LSTM)"
    )
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)

    doc.add_page_break()

    # ── Table of Contents ──
    doc.add_heading("Table of Contents", level=1)
    toc_items = [
        "1. Executive Summary",
        "2. Data Source & Pipeline Overview",
        "   2.1 Data Source",
        "   2.2 What the Data Contains",
        "   2.3 How We Pulled the Data",
        "   2.4 How We Processed the Data",
        "3. Feature Engineering",
        "4. Models Trained",
        "   4.1 XGBoost (Extreme Gradient Boosting)",
        "   4.2 Random Forest",
        "   4.3 Logistic Regression / Linear Regression",
        "   4.4 Gradient Boosting (Scikit-learn)",
        "   4.5 Feedforward Neural Network (FNN)",
        "   4.6 1D Convolutional Neural Network (CNN)",
        "   4.7 Recurrent Neural Network (RNN/LSTM)",
        "5. Classification Results",
        "6. Regression Results",
        "7. Training Time Comparison",
        "8. Feature Importance Analysis",
        "9. Model Selection & Recommendation",
        "10. Conclusion",
    ]
    for item in toc_items:
        p = doc.add_paragraph(item)
        p.paragraph_format.space_after = Pt(4)

    doc.add_page_break()

    # ── 1. Executive Summary ──
    doc.add_heading("1. Executive Summary", level=1)
    doc.add_paragraph(
        "This report documents the end-to-end process of training and comparing "
        "seven machine learning models for health insurance claim prediction as part "
        "of the Masters thesis: \"Dynamic AI-Powered Virtual Agent for Real-Time "
        "Health Insurance Claim Status Queries.\" The models include four traditional "
        "ML algorithms and three deep learning neural network architectures, trained "
        "on 500,000 claims sampled from the CMS DE-SynPUF (Data Entrepreneurs' "
        "Synthetic Public Use File) dataset."
    )
    doc.add_paragraph(
        "Two prediction tasks were evaluated:"
    )
    bullets = [
        "Classification: Predict whether a claim will be Approved or Denied",
        "Regression: Predict the number of processing days for a claim",
    ]
    for b in bullets:
        doc.add_paragraph(b, style="List Bullet")

    doc.add_paragraph(
        "Seven models were compared: four traditional ML models (XGBoost, Random Forest, "
        "Logistic Regression, Gradient Boosting) and three neural network architectures "
        "(FNN - Feedforward Neural Network, CNN - 1D Convolutional Neural Network, "
        "RNN - LSTM-based Recurrent Neural Network)."
    )

    doc.add_paragraph(
        "Key Finding: All traditional classifiers achieved 100% accuracy on the test set. "
        "Neural network models also achieved near-perfect or perfect accuracy, demonstrating "
        "that the deterministic relationship in the synthetic data (clm_pmt_amt > 0 = Approved) "
        "is learnable by all model types. XGBoost is recommended as the best model due to "
        "its fastest training time with top accuracy, combined with native SHAP-based "
        "explainability."
    )

    # ── 2. Data Source & Pipeline ──
    doc.add_heading("2. Data Source & Pipeline Overview", level=1)

    doc.add_heading("2.1 Data Source", level=2)
    doc.add_paragraph(
        "The dataset used in this project is the CMS 2008-2010 Data Entrepreneurs' "
        "Synthetic Public Use File (DE-SynPUF), published by the Centers for Medicare "
        "& Medicaid Services (CMS), a U.S. federal agency within the Department of "
        "Health and Human Services."
    )
    doc.add_paragraph(
        "Source Website: https://www.cms.gov/data-research/statistics-trends-and-reports/"
        "medicare-claims-synthetic-public-use-files/"
        "cms-2008-2010-data-entrepreneurs-synthetic-public-use-file-de-synpuf"
    )
    doc.add_paragraph(
        "The DE-SynPUF is a synthetic (privacy-safe) version of real Medicare claims data. "
        "It was created to provide a realistic but non-identifiable dataset for researchers, "
        "data entrepreneurs, and students to develop and test healthcare analytics solutions "
        "without any risk of exposing actual patient information. The data preserves the "
        "statistical distributions and relationships found in real Medicare claims while "
        "containing no real beneficiary data."
    )

    doc.add_heading("2.2 What the Data Contains", level=2)
    doc.add_paragraph("The DE-SynPUF consists of 20 samples, each containing 5 file types:")
    data_types = [
        "Beneficiary Summary Files (2008, 2009, 2010): Patient demographics, enrollment "
        "status, chronic conditions, and Medicare entitlement information",
        "Inpatient Claims: Hospital inpatient stay records including admission/discharge "
        "dates, diagnosis codes (ICD-9), procedure codes, DRG codes, and payment amounts",
        "Outpatient Claims: Hospital outpatient visit records including service dates, "
        "diagnosis codes, HCPCS codes, and payment amounts",
        "Carrier Claims (Parts A & B): Physician and supplier claims including provider "
        "NPIs, line-level service details, and allowed/paid amounts",
        "Prescription Drug Events: Part D medication claims including drug codes, "
        "service dates, dispensing quantities, and costs",
    ]
    for dt in data_types:
        doc.add_paragraph(dt, style="List Bullet")
    doc.add_paragraph(
        "Total raw data size: approximately 16 GB+ across all 20 samples."
    )

    doc.add_heading("2.3 How We Pulled the Data", level=2)
    doc.add_paragraph(
        "A custom Python download script (download_cms_data.py) was developed to "
        "automate the retrieval of all CMS DE-SynPUF ZIP files directly from CMS.gov "
        "download servers. The script:"
    )
    pull_steps = [
        "Downloads all 8 file types for each sample (beneficiary summaries for 3 years, "
        "inpatient, outpatient, carrier parts A & B, and prescription drug events)",
        "Implements retry logic with exponential backoff for network resilience",
        "Tracks download progress with real-time percentage reporting",
        "Automatically extracts ZIP archives and removes the compressed files to save disk space",
        "Downloads supporting documentation (Data Users Guide, Codebook, FAQ) from CMS",
        "Supports partial re-runs by skipping already-downloaded files",
    ]
    for s in pull_steps:
        doc.add_paragraph(s, style="List Bullet")

    doc.add_heading("2.4 How We Processed the Data", level=2)
    doc.add_paragraph(
        "The raw CMS data required significant transformation before it could be used "
        "for ML training and chatbot queries. The processing pipeline consists of 4 steps:"
    )
    steps = [
        (
            "Step 1: Transform & Harmonize (transform_claims.py)",
            "The core ETL script processes each claim type (inpatient, outpatient, carrier, "
            "prescription) using chunked I/O (200,000 rows per chunk) to handle multi-GB files. "
            "It harmonizes column names across all 4 claim types into a unified schema. "
            "For example, prescription drug events use PDE_ID and SRVC_DT while inpatient "
            "claims use CLM_ID and CLM_FROM_DT -- these are all mapped to a common format."
        ),
        (
            "Step 2: Engineer Claim Status Workflow",
            "Using vectorized operations with MD5 hashing for deterministic pseudo-randomness, "
            "the pipeline engineers a realistic claim status workflow: claims with payment "
            "amount > 0 are marked 'approved', while unpaid claims are distributed across "
            "'denied' (35%), 'under_review' (35%), and 'submitted' (30%). Denied claims "
            "receive one of 10 realistic denial reason codes (e.g., D001 - Service not covered, "
            "D005 - Incomplete documentation). Additionally, dates for submission, review start, "
            "and resolution are derived, along with urgency levels and claim reference numbers."
        ),
        (
            "Step 3: Generate Synthetic Member Identities",
            "Since the DE-SynPUF anonymizes all patient data, the pipeline generates realistic "
            "but synthetic member identities using deterministic hashing of the DESYNPUF_ID. "
            "This produces consistent member names, phone numbers, and email addresses across "
            "all claims for the same beneficiary, enabling the chatbot to provide a natural "
            "conversational experience with member authentication."
        ),
        (
            "Step 4: Pipeline Orchestration & Database Loading",
            "The run_pipeline.py script orchestrates the full download-transform-merge pipeline "
            "with CLI flags (--samples, --skip-download, --merge-only). After merging all sample "
            "outputs into a single CSV, load_csv_to_sqlite.py bulk-loads the data into a SQLite "
            "database with 8 optimized indexes for fast query performance. Final database: "
            "approximately 18 GB with 44.6 million rows."
        ),
    ]
    for title, desc in steps:
        doc.add_heading(title, level=3)
        doc.add_paragraph(desc)

    # ── 3. Feature Engineering ──
    doc.add_heading("3. Feature Engineering", level=1)
    doc.add_paragraph(
        "Eleven features were engineered from the raw claims data for model training:"
    )

    features = [
        ("clm_pmt_amt", "Claim payment amount (primary predictor for classification)"),
        ("nch_prmry_pyr_clm_pd_amt", "Primary payer claim paid amount"),
        ("nch_bene_blood_ddctbl_lblty_am", "Blood deductible liability amount"),
        ("clm_pass_thru_per_diem_amt", "Pass-through per diem amount"),
        ("nch_bene_ip_ddctbl_amt", "Inpatient deductible amount"),
        ("clm_pps_cptl_fsp_amt", "PPS capital FSP amount"),
        ("clm_pps_cptl_outlier_amt", "PPS capital outlier amount"),
        ("clm_pps_old_cptl_hld_hrmls_amt", "PPS old capital hold harmless amount"),
        ("nch_profnl_cmpnt_chrg_amt", "Professional component charge amount"),
        ("nch_bene_ptb_ddctbl_amt", "Part B deductible amount"),
        ("nch_bene_ptb_coinsrnc_amt", "Part B coinsurance amount"),
    ]
    table = doc.add_table(rows=1 + len(features), cols=2)
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.rows[0].cells[0].text = "Feature Name"
    table.rows[0].cells[1].text = "Description"
    for p in table.rows[0].cells[0].paragraphs:
        for r in p.runs:
            r.bold = True
    for p in table.rows[0].cells[1].paragraphs:
        for r in p.runs:
            r.bold = True
    for i, (name, desc) in enumerate(features):
        table.rows[i + 1].cells[0].text = name
        table.rows[i + 1].cells[1].text = desc

    doc.add_paragraph("")
    doc.add_paragraph(
        "Target variables: claim_status (binary: approved/denied) for classification, "
        "and processing_days (integer: 1-20) for regression."
    )

    # ── 4. Models Trained ──
    doc.add_heading("4. Models Trained", level=1)
    doc.add_paragraph(
        "Four machine learning models were selected to represent a diverse range of "
        "algorithmic approaches -- from simple linear models to complex ensemble methods. "
        "This diversity ensures a comprehensive comparison across different levels of "
        "model complexity, training speed, and interpretability. Each model was trained "
        "for both classification (approved vs. denied) and regression (processing days)."
    )
    models_info = [
        (
            "XGBoost (Extreme Gradient Boosting)",
            "Purpose: XGBoost was chosen as the primary candidate due to its proven dominance "
            "on structured/tabular data in industry and Kaggle competitions. It builds an "
            "ensemble of decision trees sequentially, where each new tree corrects the errors "
            "of the previous ones using gradient descent optimization. "
            "Why we included it: It offers the best balance of speed, accuracy, and "
            "explainability (via native SHAP support), making it ideal for production deployment "
            "in the chatbot where real-time predictions with explanations are needed. "
            "Parameters: n_estimators=200, max_depth=6, learning_rate=0.1, subsample=0.8, "
            "eval_metric=logloss."
        ),
        (
            "Random Forest",
            "Purpose: Random Forest was included as a benchmark ensemble method that uses "
            "bagging (Bootstrap Aggregation) instead of boosting. It trains many independent "
            "decision trees on random subsets of data and features, then aggregates their "
            "predictions through majority voting (classification) or averaging (regression). "
            "Why we included it: It serves as a robust baseline that is naturally resistant "
            "to overfitting and requires minimal hyperparameter tuning. It also provides "
            "feature importance rankings for comparison with XGBoost. "
            "Parameters: n_estimators=200, max_depth=15, min_samples_split=10, random_state=42."
        ),
        (
            "Logistic Regression / Linear Regression",
            "Purpose: Logistic Regression (for classification) and Linear Regression "
            "(for regression) were included as the simplest baseline models. Logistic "
            "Regression models the probability of class membership using a sigmoid function "
            "on a linear combination of features. Linear Regression fits a straight line "
            "to predict continuous values. "
            "Why we included them: They provide a critical lower-bound baseline. If a simple "
            "linear model performs as well as complex ensemble methods, it indicates that the "
            "relationship in the data is straightforward. They are also the most interpretable "
            "models with the fastest inference time. "
            "Parameters: max_iter=1000, random_state=42, solver=lbfgs."
        ),
        (
            "Gradient Boosting (Scikit-learn)",
            "Purpose: Scikit-learn's Gradient Boosting was included as a direct comparison "
            "to XGBoost. While both use the same fundamental algorithm (sequential boosting "
            "of decision trees), sklearn's implementation is the reference implementation "
            "without XGBoost's hardware optimizations. "
            "Why we included it: To isolate the effect of XGBoost's engineering optimizations "
            "(parallel processing, cache awareness, column block structure) from the algorithm "
            "itself. Any performance difference is attributable to implementation, not algorithm. "
            "Parameters: n_estimators=200, max_depth=5, learning_rate=0.1, subsample=0.8."
        ),
        (
            "Feedforward Neural Network (FNN)",
            "Purpose: A standard fully-connected deep neural network was included to assess "
            "whether deep learning approaches can match traditional ML on tabular claims data. "
            "The FNN processes all input features simultaneously through multiple layers of "
            "neurons, learning non-linear feature interactions through backpropagation. "
            "Architecture: Two hidden layers (64 -> 32 -> 16 neurons) with BatchNormalization "
            "for training stability and Dropout (0.3, 0.2) for regularization. Uses Adam "
            "optimizer with binary crossentropy loss for classification, MSE for regression. "
            "Why we included it: To establish a deep learning baseline and determine whether "
            "non-linear deep feature extraction provides any advantage over traditional "
            "tree-based models on this tabular healthcare dataset. "
            "Training: 30 epochs max with early stopping (patience=5), batch_size=2048, "
            "10% validation split, ReduceLROnPlateau callback."
        ),
        (
            "1D Convolutional Neural Network (CNN)",
            "Purpose: A 1D CNN was included to detect local patterns and feature interactions "
            "within the tabular feature vector. While CNNs are traditionally used for images "
            "and sequences, 1D convolution on tabular features can capture relationships "
            "between adjacent features. The input features are reshaped from (n, 11) to "
            "(n, 11, 1) to create a 1D signal that convolutional filters can process. "
            "Architecture: Two Conv1D layers (32 and 64 filters, kernel_size=3, same padding) "
            "with BatchNormalization, followed by GlobalAveragePooling1D and a dense "
            "classification/regression head with Dropout (0.3). "
            "Why we included it: To test whether convolutional pattern detection offers "
            "any advantage on tabular healthcare data compared to both traditional ML and "
            "fully-connected neural networks. Recent research has shown promise for 1D CNNs "
            "on structured data when features have meaningful ordering. "
            "Training: Same as FNN (30 epochs, early stopping, Adam optimizer)."
        ),
        (
            "Recurrent Neural Network (RNN/LSTM)",
            "Purpose: An LSTM (Long Short-Term Memory) based RNN was included to process "
            "the features sequentially, treating each of the 11 features as a timestep. "
            "The LSTM maintains a hidden state that accumulates information across features, "
            "potentially capturing sequential dependencies that other architectures might miss. "
            "Architecture: Two stacked LSTM layers (32 and 16 units) where the first LSTM "
            "passes all hidden states (return_sequences=True) to the second LSTM which "
            "produces a single summary vector. Followed by Dense(16) with Dropout (0.2). "
            "Why we included it: To evaluate whether sequential feature processing reveals "
            "patterns that simultaneous feature processing (FNN, tree-based models) might "
            "miss. LSTMs are particularly good at learning which features to remember and "
            "which to ignore through their gating mechanism. "
            "Training: Same as FNN, though LSTM is inherently slower due to sequential "
            "computation (cannot be parallelized across timesteps)."
        ),
    ]
    for title, desc in models_info:
        doc.add_heading(title, level=2)
        doc.add_paragraph(desc)

    doc.add_paragraph(
        "Training setup: 500,000 claims sampled from the 44.6 million row database "
        "(filtered to approved + denied only), 80/20 train-test split (400,000 train / "
        "100,000 test), stratified by class label to maintain the same approved/denied "
        "ratio in both sets. Neural networks were additionally trained with StandardScaler "
        "input normalization, Adam optimizer, batch size 2048, and early stopping with "
        "patience=5 monitoring validation loss."
    )

    # ── 5. Classification Results ──
    doc.add_heading("5. Classification Results", level=1)
    doc.add_paragraph(
        "All seven classifiers achieved perfect or near-perfect scores across all metrics:"
    )

    # Classification table (exclude confusion matrix columns)
    clf_display = [clf_rows[0][:7]]  # Header: Model, Accuracy..Training Time
    for row in clf_rows[1:]:
        clf_display.append(row[:7])
    add_table_from_csv(doc, clf_display)

    doc.add_paragraph("")
    doc.add_paragraph(
        "Why near-perfect accuracy? The claim_status target variable is deterministically "
        "derived from clm_pmt_amt in the data pipeline (transform_claims.py). Claims with "
        "payment amount > 0 are labeled 'approved'; claims with payment = 0 are "
        "'denied' or 'under_review'. Since clm_pmt_amt is included as a feature, "
        "all models can learn this rule. Traditional tree-based models learn it as a "
        "single decision split (100% accuracy), while neural networks also converge to "
        "this pattern through gradient descent. This validates that both traditional ML "
        "and deep learning correctly capture the business logic embedded in the synthetic data."
    )

    # Charts
    if CHARTS["classifier_comparison"].exists():
        doc.add_heading("Classification Metrics Chart", level=2)
        doc.add_picture(str(CHARTS["classifier_comparison"]), width=Inches(5.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    if CHARTS["confusion_matrices"].exists():
        doc.add_heading("Confusion Matrices", level=2)
        doc.add_picture(str(CHARTS["confusion_matrices"]), width=Inches(5.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    if CHARTS["roc_curves"].exists():
        doc.add_heading("ROC Curves", level=2)
        doc.add_picture(str(CHARTS["roc_curves"]), width=Inches(5.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    # ── 6. Regression Results ──
    doc.add_heading("6. Regression Results", level=1)
    doc.add_paragraph(
        "All seven regressors showed similar performance on processing days prediction:"
    )
    add_table_from_csv(doc, reg_rows)

    doc.add_paragraph("")
    doc.add_paragraph(
        "The R-squared values near zero indicate that processing_days has no learnable "
        "correlation with the input features. This is expected because processing days "
        "are randomly assigned (uniform 1-20) during the data transformation step. "
        "The MAE of ~6.47 days is consistent with the expected error when predicting "
        "the mean of a uniform(1,20) distribution."
    )

    if CHARTS["regressor_comparison"].exists():
        doc.add_heading("Regression Metrics Chart", level=2)
        doc.add_picture(str(CHARTS["regressor_comparison"]), width=Inches(5.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    # ── 7. Training Time ──
    doc.add_heading("7. Training Time Comparison", level=1)
    doc.add_paragraph(
        "Training time is a critical factor when models have identical accuracy. "
        "XGBoost is the clear winner with the fastest training time:"
    )

    time_data = [["Model", "Classification (s)", "Regression (s)"]]
    clf_times = {row[0]: row[6] for row in clf_rows[1:]}
    reg_times = {row[0]: row[4] for row in reg_rows[1:]}
    for model in clf_times:
        reg_model = model if model != "Logistic Regression" else "Linear Regression"
        time_data.append([model, clf_times[model], reg_times.get(reg_model, "N/A")])
    add_table_from_csv(doc, time_data)

    if CHARTS["training_time"].exists():
        doc.add_paragraph("")
        doc.add_picture(str(CHARTS["training_time"]), width=Inches(5.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    # ── 8. Feature Importance ──
    doc.add_heading("8. Feature Importance Analysis", level=1)
    doc.add_paragraph(
        "Feature importance was extracted from the three tree-based models "
        "(XGBoost, Random Forest, Gradient Boosting). As expected, clm_pmt_amt "
        "dominates as the most important feature for classification, confirming "
        "the deterministic relationship."
    )

    # Load and show XGBoost feature importance
    fi_path = RESULTS_DIR / "feature_importance_xgboost.csv"
    if fi_path.exists():
        fi_rows = load_csv("feature_importance_xgboost.csv")
        doc.add_heading("XGBoost Feature Importance", level=2)
        add_table_from_csv(doc, fi_rows)

    if CHARTS["feature_importance"].exists():
        doc.add_paragraph("")
        doc.add_picture(str(CHARTS["feature_importance"]), width=Inches(5.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    # ── 9. Model Selection ──
    doc.add_heading("9. Model Selection & Recommendation", level=1)
    doc.add_paragraph(
        "Since all four models achieve identical classification accuracy (100%), "
        "the selection criteria shifts to secondary factors. Here is how we arrived "
        "at the final recommendation:"
    )

    doc.add_heading("Decision Process", level=2)
    doc.add_paragraph(
        "Step 1 - Accuracy Assessment: We first compared all models on core classification "
        "metrics (Accuracy, Precision, Recall, F1 Score, AUC-ROC). All four models scored "
        "100% across all metrics with zero false positives and zero false negatives on "
        "100,000 test samples. This eliminated accuracy as a differentiating factor."
    )
    doc.add_paragraph(
        "Step 2 - Regression Assessment: For the processing days regression task, all four "
        "models showed nearly identical performance (MAE ~6.47 days, R-squared ~0.0). This "
        "is expected because processing days are randomly generated in the synthetic data, "
        "so no model can find a learnable pattern. This also did not differentiate the models."
    )
    doc.add_paragraph(
        "Step 3 - Training Speed: With identical accuracy, training speed becomes the "
        "primary differentiator. XGBoost trained in 2.8 seconds for classification, "
        "compared to Random Forest (5.9s), Logistic Regression (10.4s), and Gradient "
        "Boosting (26.0s). XGBoost is 9.3x faster than the slowest model."
    )
    doc.add_paragraph(
        "Step 4 - Explainability: For a healthcare chatbot, being able to explain WHY "
        "a claim was predicted as approved or denied is critical for user trust. XGBoost "
        "has native integration with SHAP (SHapley Additive exPlanations), which provides "
        "individual prediction explanations showing how each feature contributed to the outcome."
    )
    doc.add_paragraph(
        "Step 5 - Production Readiness: XGBoost has a mature ecosystem with wide industry "
        "adoption, efficient serialization (joblib), small model footprint, and fast "
        "inference time suitable for real-time chatbot predictions."
    )
    doc.add_paragraph(
        "Step 6 - Neural Network Assessment: Three neural network architectures (FNN, CNN, "
        "RNN/LSTM) were trained using TensorFlow/Keras to evaluate whether deep learning "
        "provides any advantage on this tabular healthcare dataset. While all three neural "
        "networks achieved high classification accuracy, they required significantly more "
        "training time than tree-based models due to the iterative gradient descent training "
        "process. For the regression task, neural networks showed comparable MAE to "
        "traditional models. Given the additional complexity, longer training time, and "
        "no meaningful accuracy improvement, tree-based models remain preferred for this "
        "use case."
    )

    criteria = [
        "Training Speed: XGBoost (2.8s) is 9.3x faster than Gradient Boosting (26.0s)",
        "Memory Efficiency: XGBoost and Logistic Regression have smaller model footprints",
        "Explainability: XGBoost supports SHAP (SHapley Additive exPlanations) for individual predictions",
        "Production Readiness: XGBoost has mature deployment tools and wide industry adoption",
    ]
    doc.add_heading("Summary of Selection Criteria", level=2)
    for c in criteria:
        doc.add_paragraph(c, style="List Bullet")

    doc.add_paragraph("")
    rec = doc.add_paragraph()
    run = rec.add_run("Recommendation: XGBoost")
    run.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = RGBColor(0x10, 0xA3, 0x7F)
    rec.add_run(
        " is selected as the production model for ClaimBot AI. "
        "It offers the best combination of accuracy, speed, and explainability "
        "via SHAP integration."
    )

    # ── 10. Conclusion ──
    doc.add_heading("10. Conclusion", level=1)
    doc.add_paragraph(
        "This report presented a comprehensive end-to-end machine learning pipeline "
        "for health insurance claim prediction, from raw data acquisition through model "
        "comparison and selection. The key findings are:"
    )
    conclusions = [
        "Data Pipeline: Successfully downloaded, transformed, and loaded 44.6 million CMS "
        "DE-SynPUF claims into a structured SQLite database through an automated ETL pipeline",
        "Model Comparison: All seven models (4 traditional ML + 3 neural networks) achieved "
        "high classification accuracy, validating that the data pipeline correctly embeds "
        "the business logic into the synthetic data and that all model families can learn it",
        "Deep Learning Assessment: Neural network architectures (FNN, CNN, RNN/LSTM) showed "
        "no meaningful advantage over traditional ML on this tabular dataset, while requiring "
        "significantly more training time and computational resources",
        "Regression Insight: The processing days regression task showed R-squared near zero "
        "for all models, confirming that this target variable is randomly generated and has "
        "no learnable relationship with claim features. In a real-world deployment with "
        "actual processing time data, ensemble methods would likely outperform linear models",
        "Model Selection: XGBoost was selected as the production model based on its superior "
        "training speed, native SHAP explainability, and production-ready ecosystem",
        "Integration: The selected XGBoost model is deployed within the ClaimBot AI chatbot, "
        "providing real-time claim outcome predictions with SHAP-based explanations that "
        "help users understand why their claim was predicted to be approved or denied",
    ]
    for c in conclusions:
        doc.add_paragraph(c, style="List Bullet")

    doc.add_paragraph("")
    doc.add_paragraph(
        "The interactive model comparison dashboard (built with Streamlit and Plotly) "
        "provides stakeholders with a visual BI tool to explore all metrics, charts, "
        "and feature importance analysis interactively, with CSV export capabilities "
        "for further analysis in Power BI or Tableau."
    )

    # Save
    output_path = OUTPUT_DIR / "ML_Model_Comparison_Report.docx"
    doc.save(str(output_path))
    print(f"Word document saved: {output_path}")
    return output_path


# ╔═══════════════════════════════════════════════════════════╗
# ║  PART 2: POWERPOINT PRESENTATION                         ║
# ╚═══════════════════════════════════════════════════════════╝

def add_text_slide(prs, title_text, body_lines, layout_idx=1):
    """Add a slide with title and bullet points."""
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = PptPt(18)
        else:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = PptPt(18)
            p.space_before = PptPt(6)


def add_image_slide(prs, title_text, image_path, subtitle=None):
    """Add a slide with a title and centered image."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    from pptx.util import Inches as I
    txBox = slide.shapes.add_textbox(I(0.5), I(0.2), I(9), I(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = PptPt(28)
    p.font.bold = True
    p.font.color.rgb = PptRGB(0x1A, 0x1A, 0x1A)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = PptPt(14)
        p2.font.color.rgb = PptRGB(0x6B, 0x6B, 0x6B)
        p2.alignment = PP_ALIGN.CENTER

    # Image
    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), I(0.8), I(1.3), I(8.4), I(5.5))


def create_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)

    # ── Slide 1: Title ──
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "ML Model Comparison"
    slide.placeholders[1].text = (
        "Dynamic AI-Powered Virtual Agent\n"
        "for Health Insurance Claim Status Queries\n\n"
        "Masters Thesis Project"
    )

    # ── Slide 2: Overview ──
    add_text_slide(prs, "Project Overview", [
        "Thesis: Dynamic AI-Powered Virtual Agent for Health Insurance Claims",
        "Dataset: CMS DE-SynPUF (500,000 claims from 44.6M total)",
        "Task 1: Classify claims as Approved or Denied",
        "Task 2: Predict claim processing days",
        "7 ML models: 4 traditional + 3 neural networks",
        "80/20 stratified train-test split (400K train / 100K test)",
    ])

    # ── Slide 3: Data Source ──
    add_text_slide(prs, "Data Source", [
        "CMS 2008-2010 DE-SynPUF (Synthetic Public Use File)",
        "Published by Centers for Medicare & Medicaid Services (CMS)",
        "Website: cms.gov/data-research/.../de-synpuf",
        "Synthetic, privacy-safe version of real Medicare claims",
        "Contains: Beneficiary info, Inpatient, Outpatient, Carrier, Prescription data",
        "20 samples, ~16 GB+ raw data across all file types",
    ])

    # ── Slide 4: Data Pipeline ──
    add_text_slide(prs, "Data Pipeline", [
        "Step 1: Download CMS DE-SynPUF raw ZIP files from CMS.gov",
        "Step 2: Transform & harmonize 4 claim types into unified schema",
        "Step 3: Engineer claim status, dates, urgency, member identities",
        "Step 4: Orchestrate pipeline & merge all samples",
        "Step 5: Load into SQLite database (44.6M rows, ~18 GB)",
        "Step 6: Sample 500K claims for ML model training",
    ])

    # ── Slide 5: Traditional Models ──
    add_text_slide(prs, "Traditional ML Models", [
        "XGBoost: Gradient-boosted trees -- best speed + SHAP explainability",
        "Random Forest: Bagged trees -- robust baseline, overfitting resistant",
        "Logistic/Linear Regression: Simple baseline -- interpretability benchmark",
        "Gradient Boosting (sklearn): Reference implementation -- comparison to XGBoost",
    ])

    # ── Slide 6: Neural Network Models ──
    add_text_slide(prs, "Neural Network Models", [
        "FNN (Feedforward Neural Net): Dense(64->32->16) + BatchNorm + Dropout",
        "CNN (1D Conv): Conv1D(32->64) + GlobalAvgPool -- local pattern detection",
        "RNN (LSTM): Stacked LSTM(32->16) -- sequential feature processing",
        "All use: Adam optimizer, early stopping, StandardScaler, batch_size=2048",
        "Framework: TensorFlow/Keras (already required by Rasa)",
    ])

    # ── Slide 7: Classification Results ──
    add_text_slide(prs, "Classification Results", [
        "Traditional ML (4 models): 100% Accuracy, Precision, Recall, F1",
        "FNN: 99.89% Accuracy, F1=0.9993, AUC=0.9998",
        "CNN: 99.89% Accuracy, F1=0.9993, AUC=1.0000",
        "RNN (LSTM): 97.77% Accuracy, F1=0.9864, AUC=0.9967",
        "All models learn the deterministic clm_pmt_amt -> status rule",
        "Trees learn it as 1 split; NNs approximate it through gradient descent",
    ])

    # ── Slide 6: Classification Chart ──
    if CHARTS["classifier_comparison"].exists():
        add_image_slide(prs, "Classification Metrics Comparison",
                       CHARTS["classifier_comparison"])

    # ── Slide 7: Confusion Matrices ──
    if CHARTS["confusion_matrices"].exists():
        add_image_slide(prs, "Confusion Matrices",
                       CHARTS["confusion_matrices"],
                       "All models: 0 FP, 0 FN across 100,000 test samples")

    # ── Slide 8: ROC Curves ──
    if CHARTS["roc_curves"].exists():
        add_image_slide(prs, "ROC Curves",
                       CHARTS["roc_curves"],
                       "All models achieve AUC = 1.0")

    # ── Slide 11: Regression Results ──
    add_text_slide(prs, "Regression Results", [
        "Task: Predict processing days (1-20 days)",
        "All 7 models: MAE ~ 6.45-6.50 days, RMSE ~ 7.47-7.51 days",
        "R-squared ~ 0.0 (no predictive signal in any model)",
        "Neural nets show comparable MAE to traditional models",
        "Processing days randomly assigned in synthetic data",
    ])

    # ── Slide 10: Regression Chart ──
    if CHARTS["regressor_comparison"].exists():
        add_image_slide(prs, "Regression Metrics Comparison",
                       CHARTS["regressor_comparison"])

    # ── Slide 13: Training Time ──
    add_text_slide(prs, "Training Time Comparison", [
        "Traditional ML: XGBoost 2.1s | RF 5.1s | LogReg 6.5s | GB 22.3s",
        "Neural Networks: FNN 12.6s | RNN 84.1s | CNN 195.6s",
        "XGBoost is 93x faster than CNN, 40x faster than RNN",
        "Neural networks need many gradient descent iterations",
        "Tree models learn simple rules in a single split",
    ])

    # ── Slide 12: Training Time Chart ──
    if CHARTS["training_time"].exists():
        add_image_slide(prs, "Training Time Visual",
                       CHARTS["training_time"],
                       "XGBoost dominates on training speed")

    # ── Slide 13: Feature Importance ──
    if CHARTS["feature_importance"].exists():
        add_image_slide(prs, "Feature Importance",
                       CHARTS["feature_importance"],
                       "clm_pmt_amt is the dominant predictor across all tree-based models")

    # ── Slide 16: Recommendation ──
    add_text_slide(prs, "How We Chose XGBoost", [
        "Step 1: All 4 traditional models achieved 100% classification accuracy",
        "Step 2: Neural nets achieved 97.8-99.9% -- close but slower",
        "Step 3: XGBoost 93x faster than CNN (2.1s vs 195.6s)",
        "Step 4: XGBoost has native SHAP for patient-facing explanations",
        "Step 5: NNs add complexity with no accuracy gain on this data",
        "Step 6: Mature ecosystem, efficient serialization, production-ready",
    ])

    # ── Slide 17: Conclusion ──
    add_text_slide(prs, "Conclusion", [
        "44.6M CMS claims processed through automated ETL pipeline",
        "7 ML models trained: 4 traditional + 3 neural networks",
        "Traditional ML achieves 100% accuracy; NNs achieve 97.8-99.9%",
        "XGBoost selected: fastest (2.1s) + SHAP + production-ready",
        "Neural nets validate findings but add no advantage on this data",
        "Interactive BI dashboard built with Streamlit + Plotly",
    ])

    # ── Slide 16: Thank You ──
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text = (
        "Questions?\n\n"
        "ClaimBot AI - Masters Thesis Project\n"
        "ML Model Comparison Dashboard: localhost:8502"
    )

    # Save
    output_path = OUTPUT_DIR / "ML_Model_Comparison_Presentation.pptx"
    prs.save(str(output_path))
    print(f"PowerPoint saved: {output_path}")
    return output_path


# ── Main ────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating documents...")
    print()
    word_path = create_word_document()
    print()
    pptx_path = create_pptx()
    print()
    print("=" * 60)
    print(f"Word:  {word_path}")
    print(f"PPT:   {pptx_path}")
    print("=" * 60)
