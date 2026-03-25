"""
Generate MS Word Report and PowerPoint Presentation
for Anomaly Detection - ClaimBot AI Masters Thesis
"""

import csv
from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt, Emu as PptEmu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

RESULTS_DIR = Path(__file__).resolve().parent / "results"
OUTPUT_DIR = Path(__file__).resolve().parent.parent / "documents"
OUTPUT_DIR.mkdir(exist_ok=True)

# ── Load CSV data ───────────────────────────────────────────

def load_csv(name):
    rows = []
    with open(str(RESULTS_DIR / name), "r") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(row)
    return rows


metrics_rows = load_csv("anomaly_model_metrics.csv")

# Chart image paths
CHARTS = {
    "model_comparison": RESULTS_DIR / "model_comparison.png",
    "reconstruction_error": RESULTS_DIR / "reconstruction_error_distribution.png",
    "precision_recall": RESULTS_DIR / "precision_recall_curves.png",
    "anomaly_scatter": RESULTS_DIR / "anomaly_scatter.png",
    "training_time": RESULTS_DIR / "training_time.png",
    "feature_importance": RESULTS_DIR / "feature_importance_anomaly.png",
}

GREEN = RGBColor(0x10, 0xA3, 0x7F)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)


# ── Utility: formatted table ───────────────────────────────

def add_table(doc, rows, max_cols=None):
    header = rows[0][:max_cols] if max_cols else rows[0]
    data = [r[:max_cols] if max_cols else r for r in rows[1:]]

    table = doc.add_table(rows=1 + len(data), cols=len(header))
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    for j, val in enumerate(header):
        cell = table.rows[0].cells[j]
        cell.text = val
        for p in cell.paragraphs:
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in p.runs:
                run.bold = True
                run.font.size = Pt(9)

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.rows[i + 1].cells[j]
            try:
                fval = float(val)
                cell.text = f"{fval:.4f}" if 0 < abs(fval) < 100 else f"{fval:.1f}"
            except ValueError:
                cell.text = val
            for p in cell.paragraphs:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in p.runs:
                    run.font.size = Pt(9)

    return table


def add_chart(doc, key, caption, width=6.0):
    path = CHARTS.get(key)
    if path and path.exists():
        doc.add_picture(str(path), width=Inches(width))
        last_p = doc.paragraphs[-1]
        last_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap = doc.add_paragraph(caption)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap.runs[0].font.size = Pt(9) if cap.runs else None
        cap.runs[0].font.color.rgb = GRAY if cap.runs else None


# ╔═══════════════════════════════════════════════════════════╗
# ║  PART 1: MS WORD DOCUMENT                                ║
# ╚═══════════════════════════════════════════════════════════╝

def create_word_document():
    doc = Document()

    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    style.font.color.rgb = DARK

    # ── Title Page ──
    for _ in range(6):
        doc.add_paragraph("")

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Anomaly Detection Report")
    run.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = GREEN

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(
        "Deep Learning-Based Anomaly Detection for\n"
        "Healthcare Insurance Claims Fraud Identification"
    )
    run.font.size = Pt(16)
    run.font.color.rgb = GRAY

    doc.add_paragraph("")
    info = doc.add_paragraph()
    info.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = info.add_run(
        "Masters Thesis Project\n"
        "Dataset: CMS DE-SynPUF (1,000,000 claims)\n"
        "Deep Learning: Autoencoder, VAE, CNN Autoencoder\n"
        "Traditional ML: Isolation Forest, One-Class SVM"
    )
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY

    doc.add_page_break()

    # ── Table of Contents ──
    doc.add_heading("Table of Contents", level=1)
    toc_items = [
        "1. Executive Summary",
        "2. Research Question",
        "3. Dataset & Sampling",
        "4. Feature Engineering (31 Features)",
        "   4.1 Payment Features",
        "   4.2 Temporal Features",
        "   4.3 Diagnosis Complexity Features",
        "   4.4 Provider/Physician Pattern Features",
        "   4.5 Member Behavior Features",
        "   4.6 Categorical Features",
        "5. Models Used",
        "   5.1 Autoencoder (Deep Learning)",
        "   5.2 Variational Autoencoder (Deep Learning)",
        "   5.3 1D-CNN Autoencoder (Deep Learning)",
        "   5.4 Isolation Forest (Traditional ML)",
        "   5.5 One-Class SVM (Traditional ML)",
        "6. Evaluation Methodology",
        "   6.1 Synthetic Anomaly Injection",
        "   6.2 Metrics",
        "7. Results",
        "8. Project Script Files",
        "9. Conclusion",
    ]
    for item in toc_items:
        p = doc.add_paragraph(item)
        p.paragraph_format.space_after = Pt(4)

    doc.add_page_break()

    # ── 1. Executive Summary ──
    doc.add_heading("1. Executive Summary", level=1)
    doc.add_paragraph(
        "This report presents an unsupervised anomaly detection study on health insurance "
        "claims data, conducted as part of the Masters thesis: \"Dynamic AI-Powered Virtual "
        "Agent for Real-Time Health Insurance Claim Status Queries.\" The objective is to "
        "determine whether deep learning-based autoencoders can identify potentially "
        "fraudulent or anomalous billing patterns in claims data, and how they compare "
        "to traditional anomaly detection methods."
    )
    doc.add_paragraph(
        "Five anomaly detection models were trained on 1,000,000 health insurance claims "
        "sampled from the CMS DE-SynPUF database (44.6 million total records). Three are "
        "deep learning models (Autoencoder, Variational Autoencoder, 1D-CNN Autoencoder) "
        "and two are traditional machine learning baselines (Isolation Forest, One-Class SVM)."
    )
    doc.add_paragraph(
        "Key Finding: The Autoencoder achieved the best deep learning performance with an "
        "AUC-PR of 0.984 and 100% recall at the 97.5th percentile threshold, meaning it "
        "successfully detected all injected synthetic anomalies. The One-Class SVM baseline "
        "achieved the highest overall AUC-PR of 0.988. Both outperformed the Isolation Forest "
        "(AUC-PR 0.504) and the VAE (AUC-PR 0.594), demonstrating that reconstruction-based "
        "deep learning models and distance-based methods are highly effective for claims "
        "anomaly detection."
    )

    # ── 2. Research Question ──
    doc.add_heading("2. Research Question", level=1)
    p = doc.add_paragraph()
    run = p.add_run(
        "\"Can deep learning-based anomaly detection identify potentially fraudulent or "
        "anomalous billing patterns in health insurance claims data, and how do autoencoder "
        "architectures compare to traditional anomaly detection methods?\""
    )
    run.italic = True
    run.font.size = Pt(12)

    doc.add_paragraph(
        "This question is addressed through an unsupervised learning approach where models "
        "learn the distribution of \"normal\" claims and flag deviations. Since the CMS "
        "DE-SynPUF data has no ground-truth fraud labels, we inject synthetic anomalies "
        "modeled after five known healthcare fraud patterns to enable quantitative evaluation."
    )

    # ── 3. Dataset & Sampling ──
    doc.add_heading("3. Dataset & Sampling", level=1)
    doc.add_paragraph(
        "The dataset is the CMS 2008-2010 Data Entrepreneurs' Synthetic Public Use File "
        "(DE-SynPUF), containing 44,627,170 health insurance claims across four categories: "
        "carrier (18.9M), prescription (22.2M), outpatient (3.2M), and inpatient (266K). "
        "The database is stored in SQLite format (18.6 GB) with 32 columns per claim."
    )
    doc.add_paragraph("Sampling strategy:")
    sampling_items = [
        "1,000,000 claims randomly sampled from the full 44.6M population",
        "Training set: 720,000 claims (72%) — used to train models on normal patterns",
        "Validation set: 80,000 claims (8%) — for early stopping and threshold calibration",
        "Test set: 200,000 claims (20%) — for final evaluation with synthetic anomaly injection",
    ]
    for s in sampling_items:
        doc.add_paragraph(s, style="List Bullet")

    # ── 4. Feature Engineering ──
    doc.add_heading("4. Feature Engineering (31 Features)", level=1)
    doc.add_paragraph(
        "A total of 31 numeric features were engineered from all 32 database columns, "
        "organized into six groups. This is significantly richer than the 11 features used "
        "in the original supervised model comparison, enabling more nuanced anomaly detection."
    )

    doc.add_heading("4.1 Payment Features (5 features)", level=2)
    payment_features = [
        "clm_pmt_amt_filled — Claim payment amount (NULL filled with 0)",
        "approved_amount_filled — Approved amount (NULL filled with 0)",
        "payment_approved_ratio — Ratio of payment to approved amount; far from 1.0 = suspicious",
        "payment_approved_diff — Absolute difference between payment and approved amount",
        "is_negative_payment — Binary flag for negative payments (10,685 in full dataset)",
    ]
    for f in payment_features:
        doc.add_paragraph(f, style="List Bullet")

    doc.add_heading("4.2 Temporal Features (8 features)", level=2)
    temporal_features = [
        "claim_duration_days — Days between service start and end dates",
        "submission_to_service_days — Gap between service and submission (late filing = suspicious)",
        "review_duration_days — Days in review (very fast or very slow = flag)",
        "submission_to_resolution_days — Total claim lifecycle duration",
        "processing_days — Estimated processing days",
        "claim_month_sin, claim_month_cos — Cyclical encoding of service month",
        "claim_day_of_week — Day of week of submission",
    ]
    for f in temporal_features:
        doc.add_paragraph(f, style="List Bullet")

    doc.add_heading("4.3 Diagnosis Complexity Features (4 features)", level=2)
    diagnosis_features = [
        "diagnosis_count — Count of non-null ICD-9 codes (0-10 per claim)",
        "primary_diag_frequency — Frequency encoding of primary diagnosis code (13,700 unique)",
        "has_secondary_diag — Binary flag for secondary diagnosis presence",
        "diag_code_rarity_score — Average inverse-frequency of first 5 diagnosis codes; "
        "clusters of rare codes = more suspicious",
    ]
    for f in diagnosis_features:
        doc.add_paragraph(f, style="List Bullet")

    doc.add_heading("4.4 Provider/Physician Pattern Features (5 features)", level=2)
    provider_features = [
        "provider_claim_volume — Claims per provider (7,040 unique providers)",
        "physician_claim_volume — Claims per physician (171,967 unique)",
        "provider_avg_payment — Mean payment for that provider",
        "provider_category_entropy — Shannon entropy of claim categories per provider",
        "physician_provider_ratio — Number of providers a physician bills through",
    ]
    for f in provider_features:
        doc.add_paragraph(f, style="List Bullet")

    doc.add_heading("4.5 Member Behavior Features (3 features)", level=2)
    member_features = [
        "member_claim_frequency — Claims per member (avg 107, max 415)",
        "member_avg_payment — Mean payment amount for this member",
        "member_unique_providers — Distinct providers visited (\"doctor shopping\" pattern)",
    ]
    for f in member_features:
        doc.add_paragraph(f, style="List Bullet")

    doc.add_heading("4.6 Categorical Features (6 features)", level=2)
    cat_features = [
        "urgency_ordinal — Ordinal encoding (routine=0 to critical=4)",
        "appeal_flag_binary — Binary appeal flag (3.3M appeals in full dataset)",
        "cat_carrier, cat_inpatient, cat_outpatient, cat_prescription — One-hot encoded claim category",
    ]
    for f in cat_features:
        doc.add_paragraph(f, style="List Bullet")

    doc.add_page_break()

    # ── 5. Models Used ──
    doc.add_heading("5. Models Used", level=1)
    doc.add_paragraph(
        "Five anomaly detection models were evaluated: three deep learning architectures "
        "and two traditional machine learning baselines. All models operate in an unsupervised "
        "setting — they learn what \"normal\" claims look like and flag deviations."
    )

    doc.add_heading("5.1 Autoencoder (Deep Learning)", level=2)
    doc.add_paragraph(
        "A symmetric dense autoencoder that compresses input features into a low-dimensional "
        "bottleneck and reconstructs them. Claims that reconstruct poorly (high MSE) are "
        "flagged as anomalous."
    )
    arch_ae = [
        "Encoder: 31 → 64 → BatchNorm → Dropout(0.2) → 32 → BatchNorm → Dropout(0.2) → 16 → 8 (bottleneck)",
        "Decoder: 8 → 16 → BatchNorm → Dropout(0.2) → 32 → BatchNorm → 64 → 31 (reconstruction)",
        "Loss: Mean Squared Error (MSE); Optimizer: Adam",
        "Training: 50 epochs, batch size 2048, early stopping (patience=7)",
    ]
    for a in arch_ae:
        doc.add_paragraph(a, style="List Bullet")

    doc.add_heading("5.2 Variational Autoencoder / VAE (Deep Learning)", level=2)
    doc.add_paragraph(
        "A probabilistic variant of the autoencoder that learns a latent distribution "
        "rather than a fixed encoding. Uses the reparameterization trick to sample from "
        "the learned latent space. The loss combines reconstruction error and KL divergence."
    )
    arch_vae = [
        "Encoder: 31 → 64 → BatchNorm → 32 → BatchNorm → (z_mean, z_log_var) → Sampling → 8 (latent dim)",
        "Decoder: 8 → 32 → BatchNorm → 64 → 31",
        "Loss: MSE + 0.001 × KL divergence (beta is low to prioritize reconstruction quality)",
        "Training: 50 epochs, batch size 2048, early stopping (patience=7)",
    ]
    for a in arch_vae:
        doc.add_paragraph(a, style="List Bullet")

    doc.add_heading("5.3 1D-CNN Autoencoder (Deep Learning)", level=2)
    doc.add_paragraph(
        "A convolutional autoencoder that treats the feature vector as a 1D signal. "
        "Convolutional layers capture local feature interactions that dense layers may miss."
    )
    arch_cnn = [
        "Input reshaped: (N, 31) → (N, 31, 1)",
        "Encoder: Conv1D(32, k=3) → BatchNorm → Conv1D(16, k=3) → BatchNorm → GlobalAvgPool → Dense(8)",
        "Decoder: Dense(31×16) → Reshape → Conv1DTranspose(16, k=3) → BatchNorm → Conv1DTranspose(32, k=3) → Conv1D(1, k=3)",
        "Loss: MSE; Optimizer: Adam; Training: 50 epochs, batch size 2048",
    ]
    for a in arch_cnn:
        doc.add_paragraph(a, style="List Bullet")

    doc.add_heading("5.4 Isolation Forest (Traditional ML)", level=2)
    doc.add_paragraph(
        "An ensemble method that isolates anomalies by randomly partitioning features. "
        "Anomalous points require fewer splits to isolate, producing shorter average path "
        "lengths in random trees. Configuration: 200 trees, contamination=0.02, max_samples=50%."
    )

    doc.add_heading("5.5 One-Class SVM (Traditional ML)", level=2)
    doc.add_paragraph(
        "A support vector method that learns a decision boundary around normal data in "
        "high-dimensional feature space using the RBF kernel. Points outside the boundary "
        "are classified as anomalies. Configuration: kernel=RBF, gamma=scale, nu=0.02. "
        "Subsampled to 50,000 training rows due to memory constraints (O(n²) complexity)."
    )

    doc.add_page_break()

    # ── 6. Evaluation Methodology ──
    doc.add_heading("6. Evaluation Methodology", level=1)

    doc.add_heading("6.1 Synthetic Anomaly Injection", level=2)
    doc.add_paragraph(
        "Since the CMS DE-SynPUF data contains no ground-truth fraud labels, synthetic "
        "anomalies were injected into the test set to enable quantitative evaluation. "
        "Five anomaly types, modeled after documented healthcare fraud patterns, were "
        "created by perturbing features of randomly selected test claims (5% of test set = 10,000 anomalies):"
    )
    anomaly_types = [
        "Extreme Payment (25%): Payment amounts multiplied by 5-20× with added offset, "
        "simulating overbilling",
        "Payment-Approval Mismatch (20%): Large discrepancies between payment and approved "
        "amounts, simulating unauthorized billing",
        "Impossible Duration (15%): Inpatient claims with 0-day duration or carrier claims "
        "spanning 300+ days, simulating fabricated services",
        "Diagnosis Stuffing (20%): Maximum diagnosis codes (10) with rare code combinations, "
        "simulating upcoding for higher reimbursement",
        "Provider Anomaly (20%): New/low-volume providers billing extremely high amounts, "
        "simulating phantom billing or kickback schemes",
    ]
    for a in anomaly_types:
        doc.add_paragraph(a, style="List Bullet")

    doc.add_heading("6.2 Metrics", level=2)
    doc.add_paragraph("Models were evaluated using:")
    eval_metrics = [
        "AUC-PR (Area Under Precision-Recall Curve): The primary ranking metric, measuring "
        "overall detection quality across all thresholds",
        "Precision: Of all flagged claims, what fraction are truly anomalous",
        "Recall: Of all anomalies, what fraction were detected",
        "F1 Score: Harmonic mean of precision and recall",
        "Anomaly threshold: Set at the 97.5th percentile of validation set scores",
    ]
    for m in eval_metrics:
        doc.add_paragraph(m, style="List Bullet")

    doc.add_page_break()

    # ── 7. Results ──
    doc.add_heading("7. Results", level=1)
    doc.add_paragraph("Model comparison at the 97.5th percentile threshold:")
    add_table(doc, metrics_rows)

    doc.add_paragraph("")
    add_chart(doc, "model_comparison", "Figure 1: Detection Metrics Comparison Across All Models")
    add_chart(doc, "precision_recall", "Figure 2: Precision-Recall Curves (Synthetic Anomaly Detection)")
    add_chart(doc, "reconstruction_error", "Figure 3: Reconstruction Error Distributions (Deep Learning Models)")
    add_chart(doc, "anomaly_scatter", "Figure 4: PCA Projection of Test Claims Colored by Anomaly Score")
    add_chart(doc, "training_time", "Figure 5: Training Time Comparison")
    add_chart(doc, "feature_importance", "Figure 6: Top Features Distinguishing Anomalous Claims")

    doc.add_paragraph("")
    doc.add_paragraph("Key observations:")
    observations = [
        "The Autoencoder (AUC-PR=0.984) and One-Class SVM (AUC-PR=0.988) both achieved "
        "100% recall, detecting every single synthetic anomaly at the 97.5% threshold.",
        "The CNN Autoencoder ranked third (AUC-PR=0.929), with slightly lower recall (90.4%) "
        "but still strong overall performance.",
        "The VAE (AUC-PR=0.594) underperformed, likely due to the low beta parameter causing "
        "the KL divergence to underconstrain the latent space, resulting in poorer reconstruction quality.",
        "Isolation Forest (AUC-PR=0.504) performed near random for this dataset, suggesting "
        "that the anomaly patterns are better captured by reconstruction-based or boundary-based methods.",
        "Training time varied significantly: One-Class SVM was fastest (5.4s on 50K subsample), "
        "while CNN Autoencoder was slowest (135.9s on 720K samples).",
    ]
    for o in observations:
        doc.add_paragraph(o, style="List Bullet")

    doc.add_page_break()

    # ── 8. Project Script Files ──
    doc.add_heading("8. Project Script Files", level=1)
    doc.add_paragraph(
        "The anomaly detection project is organized as a Python package under "
        "anomaly_detection/. Below is each file and its role in the pipeline:"
    )

    scripts = [
        ("anomaly_detection/config.py",
         "Centralized configuration file containing all paths (database, results directory), "
         "hyperparameters (encoding dimension=8, latent dimension=8, epochs=50, batch size=2048, "
         "patience=7), anomaly detection thresholds (97.5th percentile), contamination rate "
         "(0.02), and shared color palette."),

        ("anomaly_detection/data_loader.py",
         "Loads a random sample of 1,000,000 claims from the 44.6M-row SQLite database "
         "(claims.db, 18.6 GB). Uses PRAGMA cache_size for performance. Caches the sample "
         "as a Parquet file for instant reloads on subsequent runs."),

        ("anomaly_detection/feature_engineering.py",
         "Transforms raw claims into 31 numeric features across 6 groups: Payment (5), "
         "Temporal (8), Diagnosis (4), Provider (5), Member (3), Categorical (6). Handles "
         "date parsing (mixed YYYYMMDD and YYYY-MM-DD formats), NULL imputation, frequency "
         "encoding, cyclical time encoding, and Shannon entropy computation."),

        ("anomaly_detection/models/autoencoder.py",
         "Builds the vanilla Dense Autoencoder using Keras. Symmetric architecture: "
         "31→64→32→16→8→16→32→64→31 with BatchNormalization and Dropout(0.2). "
         "Compiled with Adam optimizer and MSE loss."),

        ("anomaly_detection/models/vae.py",
         "Builds the Variational Autoencoder using Keras 3 API. Includes a custom Sampling "
         "layer (reparameterization trick) and a VAELossLayer that computes reconstruction "
         "MSE + beta×KL divergence (beta=0.001). Uses keras.ops for Keras 3 compatibility."),

        ("anomaly_detection/models/cnn_autoencoder.py",
         "Builds the 1D-CNN Autoencoder. Encoder uses Conv1D→BatchNorm→Conv1D→GlobalAvgPool→Dense. "
         "Decoder uses Dense→Reshape→Conv1DTranspose layers. Input reshaped to (N, 31, 1)."),

        ("anomaly_detection/models/traditional.py",
         "Returns a dictionary of traditional anomaly detection models: Isolation Forest "
         "(200 trees, contamination=0.02) and One-Class SVM (RBF kernel, nu=0.02)."),

        ("anomaly_detection/synthetic_anomalies.py",
         "Injects 5 types of synthetic anomalies into the test set: Extreme Payment (25%), "
         "Payment-Approval Mismatch (20%), Impossible Duration (15%), Diagnosis Stuffing (20%), "
         "and Provider Anomaly (20%). Returns contaminated data with ground-truth labels."),

        ("anomaly_detection/evaluate.py",
         "Evaluation module that computes reconstruction errors (MSE per sample), anomaly "
         "scores for traditional models (inverted decision_function), thresholds at configurable "
         "percentiles, and full precision/recall/F1/AUC-PR metrics. Handles NaN values robustly."),

        ("anomaly_detection/train_anomaly_models.py",
         "Main orchestrator script. Runs the full pipeline in 8 steps: (1) load data, "
         "(2) engineer features, (3) scale and split, (4) train 3 DL models, (5) train "
         "2 traditional models, (6) inject synthetic anomalies and evaluate, (7) save results "
         "(CSVs, models, errors), (8) generate 6 PNG charts. Total runtime: ~6 minutes."),

        ("anomaly_detection/anomaly_dashboard.py",
         "Interactive Streamlit dashboard with 6 tabs: Overview (model comparison bar chart, "
         "key metrics), Reconstruction Errors (histograms with adjustable threshold slider), "
         "Anomaly Explorer (PCA scatter, filterable claims table), Model Comparison (PR curves, "
         "AUC-PR ranking), Feature Analysis (normal vs anomalous distributions), and Anomaly "
         "Patterns (category breakdown, provider rankings, detection heatmap)."),
    ]

    for filename, description in scripts:
        p = doc.add_paragraph()
        run = p.add_run(filename)
        run.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = GREEN
        doc.add_paragraph(description)
        doc.add_paragraph("")

    doc.add_heading("How to Run", level=2)
    doc.add_paragraph("Train all models and generate results:")
    p = doc.add_paragraph("python -m anomaly_detection.train_anomaly_models")
    p.runs[0].font.name = "Consolas"
    p.runs[0].font.size = Pt(10)
    doc.add_paragraph("Launch the interactive dashboard:")
    p = doc.add_paragraph("streamlit run anomaly_detection/anomaly_dashboard.py")
    p.runs[0].font.name = "Consolas"
    p.runs[0].font.size = Pt(10)

    doc.add_page_break()

    # ── 9. Conclusion ──
    doc.add_heading("9. Conclusion", level=1)
    doc.add_paragraph(
        "This study demonstrates that deep learning-based anomaly detection, particularly "
        "autoencoders, can effectively identify anomalous billing patterns in healthcare "
        "claims data. The Autoencoder achieved an AUC-PR of 0.984 with 100% recall on "
        "synthetic fraud patterns, performing comparably to the One-Class SVM baseline "
        "(AUC-PR 0.988) while offering the advantage of learning complex non-linear "
        "patterns in high-dimensional feature spaces."
    )
    doc.add_paragraph(
        "The rich feature engineering (31 features from 32 columns) proved critical — "
        "payment ratio features, provider volume patterns, and diagnosis rarity scores "
        "were the most discriminative features for separating normal from anomalous claims. "
        "The five synthetic anomaly types (overbilling, payment mismatch, impossible duration, "
        "diagnosis stuffing, and phantom providers) represent documented real-world fraud "
        "patterns and provide a practical evaluation framework."
    )
    doc.add_paragraph(
        "Limitations: (1) Evaluation uses synthetic anomalies, not real fraud labels. "
        "(2) The VAE's poor performance may be improvable with beta tuning. "
        "(3) The One-Class SVM was trained on a subsample (50K) due to O(n²) memory scaling. "
        "Future work could explore temporal sequence modeling of per-patient claim histories "
        "and integration of the anomaly detection module into the ClaimBot AI chatbot for "
        "real-time fraud flagging."
    )

    # Save
    output_path = OUTPUT_DIR / "Anomaly_Detection_Report.docx"
    doc.save(str(output_path))
    print(f"Word document saved: {output_path}")
    return output_path


# ╔═══════════════════════════════════════════════════════════╗
# ║  PART 2: POWERPOINT PRESENTATION                         ║
# ╚═══════════════════════════════════════════════════════════╝

def create_powerpoint():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    def add_bg(slide, r=0xFF, g=0xFF, b=0xFF):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PptRGB(r, g, b)

    def add_textbox(slide, left, top, width, height, text, size=18,
                    bold=False, color=None, align=PP_ALIGN.LEFT):
        tx = slide.shapes.add_textbox(PptInches(left), PptInches(top),
                                       PptInches(width), PptInches(height))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptPt(size)
        p.font.bold = bold
        p.alignment = align
        if color:
            p.font.color.rgb = color
        return tf

    def add_bullet_frame(slide, left, top, width, height, items, size=16, color=None):
        tx = slide.shapes.add_textbox(PptInches(left), PptInches(top),
                                       PptInches(width), PptInches(height))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = PptPt(size)
            p.space_after = PptPt(6)
            p.level = 0
            if color:
                p.font.color.rgb = color
        return tf

    def add_image(slide, key, left, top, width):
        path = CHARTS.get(key)
        if path and path.exists():
            slide.shapes.add_picture(str(path), PptInches(left), PptInches(top),
                                      PptInches(width))

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0x10, 0xA3, 0x7F)
    add_textbox(slide, 0.8, 1.8, 11.7, 1.5,
                "Anomaly Detection in Healthcare Claims",
                size=36, bold=True, color=PptRGB(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, 3.5, 11.7, 1.0,
                "Deep Learning vs Traditional ML for Fraud Identification",
                size=22, color=PptRGB(0xE0, 0xFF, 0xF0), align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, 5.0, 11.7, 1.0,
                "Masters Thesis  |  CMS DE-SynPUF  |  1,000,000 Claims  |  5 Models",
                size=16, color=PptRGB(0xC0, 0xE8, 0xD8), align=PP_ALIGN.CENTER)

    # ── Slide 2: Research Question ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Research Question",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    add_textbox(slide, 1.0, 1.4, 11.3, 1.5,
                "\"Can deep learning-based anomaly detection identify potentially "
                "fraudulent or anomalous billing patterns in health insurance claims "
                "data, and how do autoencoder architectures compare to traditional "
                "anomaly detection methods?\"",
                size=20, color=PptRGB(0x33, 0x33, 0x33))
    add_bullet_frame(slide, 1.0, 3.5, 11.3, 3.5, [
        "Unsupervised approach: no fraud labels needed",
        "Models learn \"normal\" claim patterns, flag deviations",
        "5 synthetic fraud types injected for evaluation",
        "31 features engineered from 32 database columns",
        "Dataset: 1M sample from 44.6M CMS DE-SynPUF claims",
    ], size=18, color=PptRGB(0x4A, 0x4A, 0x4A))

    # ── Slide 3: Models Used ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Models Used",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    add_textbox(slide, 0.8, 1.2, 5.5, 0.5, "Deep Learning (3)",
                size=22, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    add_bullet_frame(slide, 0.8, 1.8, 5.5, 4.5, [
        "Autoencoder: 31→64→32→16→8→16→32→64→31",
        "  Dense symmetric encoder-decoder, MSE loss",
        "VAE: Variational Autoencoder",
        "  Reparameterization trick, KL divergence",
        "CNN Autoencoder: 1D Convolutional",
        "  Conv1D→BatchNorm→GlobalAvgPool→Dense",
    ], size=15, color=PptRGB(0x4A, 0x4A, 0x4A))
    add_textbox(slide, 7.0, 1.2, 5.5, 0.5, "Traditional ML (2)",
                size=22, bold=True, color=PptRGB(0x3B, 0x82, 0xF6))
    add_bullet_frame(slide, 7.0, 1.8, 5.5, 4.5, [
        "Isolation Forest: 200 trees",
        "  Random partitioning, shorter path = anomaly",
        "One-Class SVM: RBF kernel",
        "  Decision boundary around normal data",
        "  Subsampled to 50K (memory constraint)",
    ], size=15, color=PptRGB(0x4A, 0x4A, 0x4A))

    # ── Slide 4: Feature Engineering ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Feature Engineering: 31 Features from 32 Columns",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    groups = [
        ("Payment (5)", "pmt amount, approved amount,\nratio, diff, negative flag"),
        ("Temporal (8)", "claim duration, submission gap,\nreview time, cycle encoding"),
        ("Diagnosis (4)", "code count, frequency,\nsecondary flag, rarity score"),
        ("Provider (5)", "volume, avg payment,\ncategory entropy, physician ratio"),
        ("Member (3)", "claim frequency,\navg payment, unique providers"),
        ("Categorical (6)", "urgency ordinal, appeal flag,\nclaim category one-hot"),
    ]
    for i, (title, desc) in enumerate(groups):
        col = i % 3
        row = i // 3
        left = 0.8 + col * 4.0
        top = 1.4 + row * 2.8
        add_textbox(slide, left, top, 3.6, 0.5, title,
                    size=18, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
        add_textbox(slide, left, top + 0.5, 3.6, 1.8, desc,
                    size=14, color=PptRGB(0x4A, 0x4A, 0x4A))

    # ── Slide 5: Results ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Results: Model Comparison",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    add_image(slide, "model_comparison", 0.5, 1.2, 12.3)

    # ── Slide 6: PR Curves ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Precision-Recall Curves",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    add_image(slide, "precision_recall", 1.5, 1.0, 10.0)

    # ── Slide 7: Reconstruction Errors ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Reconstruction Error Distributions",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    add_image(slide, "reconstruction_error", 0.5, 1.0, 12.3)

    # ── Slide 8: PCA Scatter ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Anomaly Score Visualization (PCA)",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    add_image(slide, "anomaly_scatter", 2.0, 1.0, 9.0)

    # ── Slide 9: Feature Importance ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Top Features for Anomaly Detection",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    add_image(slide, "feature_importance", 1.5, 1.0, 10.0)

    # ── Slide 10: Script Files ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Project Script Files",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    script_items = [
        "config.py — Paths, hyperparameters, constants",
        "data_loader.py — SQLite sampling + Parquet caching",
        "feature_engineering.py — 31 features from 32 columns",
        "models/autoencoder.py — Dense Autoencoder (Keras)",
        "models/vae.py — Variational Autoencoder (Keras 3)",
        "models/cnn_autoencoder.py — 1D-CNN Autoencoder",
        "models/traditional.py — Isolation Forest + One-Class SVM",
        "synthetic_anomalies.py — 5 fraud patterns for evaluation",
        "evaluate.py — Scoring, thresholds, PR/F1/AUC-PR metrics",
        "train_anomaly_models.py — Main pipeline orchestrator",
        "anomaly_dashboard.py — 6-tab Streamlit dashboard",
    ]
    add_bullet_frame(slide, 0.8, 1.2, 11.7, 5.5, script_items,
                     size=16, color=PptRGB(0x4A, 0x4A, 0x4A))

    # ── Slide 11: Key Findings ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Key Findings",
                size=28, bold=True, color=PptRGB(0x10, 0xA3, 0x7F))
    findings = [
        "Autoencoder: AUC-PR = 0.984, 100% recall — best deep learning model",
        "One-Class SVM: AUC-PR = 0.988, 100% recall — best overall",
        "CNN Autoencoder: AUC-PR = 0.929, 90.4% recall — strong third",
        "VAE (0.594) and Isolation Forest (0.504) underperformed",
        "Payment ratio and provider volume were most discriminative features",
        "Deep learning matches or exceeds traditional ML for anomaly detection",
    ]
    add_bullet_frame(slide, 0.8, 1.3, 11.7, 5.0, findings,
                     size=20, color=PptRGB(0x33, 0x33, 0x33))

    # ── Slide 12: Conclusion ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0x10, 0xA3, 0x7F)
    add_textbox(slide, 0.8, 1.5, 11.7, 1.5,
                "Conclusion",
                size=36, bold=True, color=PptRGB(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 3.2, 10.3, 3.5,
                "Deep learning autoencoders effectively detect anomalous billing patterns "
                "in healthcare claims without requiring labeled fraud data.\n\n"
                "The Autoencoder (AUC-PR 0.984) performs comparably to the One-Class SVM "
                "(0.988), proving that unsupervised deep learning is a viable approach for "
                "healthcare fraud detection at scale.\n\n"
                "All 31 engineered features, 5 models, and interactive dashboard are available "
                "for reproducible research.",
                size=18, color=PptRGB(0xE0, 0xFF, 0xF0), align=PP_ALIGN.CENTER)

    # Save
    output_path = OUTPUT_DIR / "Anomaly_Detection_Presentation.pptx"
    prs.save(str(output_path))
    print(f"PowerPoint saved: {output_path}")
    return output_path


# ╔═══════════════════════════════════════════════════════════╗
# ║  MAIN                                                     ║
# ╚═══════════════════════════════════════════════════════════╝

if __name__ == "__main__":
    print("Generating Anomaly Detection documents...")
    create_word_document()
    create_powerpoint()
    print("Done!")
